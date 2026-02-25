"""
EduSense - Document Intelligence Service

Handles document upload, text extraction, and storage for the document intelligence system.
Extracts text from various file formats and saves to MongoDB as StudyMaterial documents.
"""

import io
from pathlib import Path
from typing import Optional
from fastapi import HTTPException, UploadFile


def extract_text_from_pdf(file_bytes: bytes) -> str:
    """
    Extract text from a PDF file using pdfplumber.
    
    Args:
        file_bytes: Raw bytes of the PDF file
        
    Returns:
        Extracted text with page markers
    """
    import pdfplumber
    
    text_parts = []
    
    try:
        with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
            for i, page in enumerate(pdf.pages):
                page_text = page.extract_text()
                if page_text and page_text.strip():
                    text_parts.append(f"--- Page {i + 1} ---\n{page_text.strip()}")
    except Exception as e:
        raise ValueError(f"Error extracting text from PDF: {str(e)}")
    
    return "\n\n".join(text_parts) if text_parts else "No readable text found in PDF."


def extract_text_from_ppt(file_bytes: bytes) -> str:
    """
    Extract text from a PowerPoint file using python-pptx.
    
    Args:
        file_bytes: Raw bytes of the PPTX file
        
    Returns:
        Extracted text with slide markers
    """
    from pptx import Presentation
    
    text_parts = []
    
    try:
        prs = Presentation(io.BytesIO(file_bytes))
        
        for i, slide in enumerate(prs.slides):
            slide_texts = []
            
            # Extract text from all shapes in the slide
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text.strip())
            
            if slide_texts:
                text_parts.append(f"--- Slide {i + 1} ---\n" + "\n".join(slide_texts))
    except Exception as e:
        raise ValueError(f"Error extracting text from PowerPoint: {str(e)}")
    
    return "\n\n".join(text_parts) if text_parts else "No readable text found in presentation."


def extract_text_from_docx(file_bytes: bytes) -> str:
    """
    Extract text from a Word document using python-docx.
    
    Args:
        file_bytes: Raw bytes of the DOCX file
        
    Returns:
        Extracted text from all paragraphs
    """
    from docx import Document
    
    text_parts = []
    
    try:
        doc = Document(io.BytesIO(file_bytes))
        
        # Extract text from all paragraphs
        for para in doc.paragraphs:
            if para.text.strip():
                text_parts.append(para.text.strip())
        
        # Also extract text from tables
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    if cell.text.strip():
                        text_parts.append(cell.text.strip())
    except Exception as e:
        raise ValueError(f"Error extracting text from Word document: {str(e)}")
    
    return "\n\n".join(text_parts) if text_parts else "No readable text found in document."


async def process_uploaded_document(file: UploadFile, user_id: str, subject: str) -> str:
    """
    Process an uploaded document: detect file type, extract text, and save to database.
    
    Args:
        file: The uploaded file from FastAPI
        user_id: ID of the user uploading the document
        subject: Subject/course this material belongs to
        
    Returns:
        The ID of the created StudyMaterial document
        
    Raises:
        HTTPException: If file type is unsupported or processing fails
    """
    from app.models.study_material import StudyMaterial
    
    # Get file extension to detect type
    file_ext = Path(file.filename).suffix.lower()
    
    # Define supported file types and their extraction functions
    extractors = {
        ".pdf": extract_text_from_pdf,
        ".pptx": extract_text_from_ppt,
        ".ppt": extract_text_from_ppt,
        ".docx": extract_text_from_docx,
        ".doc": extract_text_from_docx,
    }
    
    # Check if file type is supported
    if file_ext not in extractors:
        raise HTTPException(
            status_code=400,
            detail=f"Unsupported file type: {file_ext}. Supported types: .pdf, .pptx, .ppt, .docx, .doc"
        )
    
    try:
        # Read file bytes
        file_bytes = await file.read()
        
        # Extract text using the appropriate extractor
        extractor = extractors[file_ext]
        extracted_text = extractor(file_bytes)
        
        # Limit text length to prevent database bloat (max ~50K chars)
        if len(extracted_text) > 50000:
            extracted_text = extracted_text[:50000] + "\n\n[Content truncated due to length...]"
        
        # Create and save StudyMaterial document
        study_material = StudyMaterial(
            user_id=user_id,
            subject=subject,
            filename=file.filename,
            extracted_text=extracted_text
        )
        
        await study_material.insert()
        
        # Return the material ID as a string
        return str(study_material.id)
        
    except ValueError as ve:
        # ValueError is raised by extraction functions with specific error messages
        raise HTTPException(status_code=400, detail=str(ve))
    except Exception as e:
        # Catch-all for unexpected errors
        raise HTTPException(
            status_code=500,
            detail=f"Error processing document: {str(e)}"
        )


async def get_user_materials(user_id: str, limit: int = 50) -> list:
    """
    Retrieve all study materials for a specific user.
    
    Args:
        user_id: ID of the user
        limit: Maximum number of materials to return (default: 50)
        
    Returns:
        List of StudyMaterial documents
    """
    from app.models.study_material import StudyMaterial
    
    materials = await StudyMaterial.find(
        StudyMaterial.user_id == user_id
    ).sort("-created_at").limit(limit).to_list()
    
    return materials


async def get_material_by_id(material_id: str, user_id: str) -> Optional[object]:
    """
    Retrieve a specific study material by ID, ensuring it belongs to the user.
    
    Args:
        material_id: ID of the study material
        user_id: ID of the user (for authorization)
        
    Returns:
        StudyMaterial document or None if not found/unauthorized
    """
    from app.models.study_material import StudyMaterial
    from beanie import PydanticObjectId
    
    try:
        material = await StudyMaterial.get(PydanticObjectId(material_id))
        
        # Verify the material belongs to the requesting user
        if material and material.user_id == user_id:
            return material
        return None
    except:
        return None
