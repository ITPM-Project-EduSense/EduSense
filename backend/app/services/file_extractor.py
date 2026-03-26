"""
EduSense - File Text Extraction Service

Extracts text content from uploaded files:
- PDF (.pdf)
- PowerPoint (.pptx)
- Word (.docx)
- Images (.png, .jpg, .jpeg) - returns description for AI
"""

import io
from pathlib import Path
from typing import List


def extract_text_from_pdf(file_bytes: bytes) -> str:
    """Extract text from a PDF file."""
    from PyPDF2 import PdfReader

    reader = PdfReader(io.BytesIO(file_bytes))
    text_parts = []

    for i, page in enumerate(reader.pages):
        page_text = page.extract_text()
        if page_text and page_text.strip():
            text_parts.append(f"--- Page {i + 1} ---\n{page_text.strip()}")

    return "\n\n".join(text_parts) if text_parts else "No readable text found in PDF."


def extract_text_from_pptx(file_bytes: bytes) -> str:
    """Extract text from a PowerPoint file."""
    from pptx import Presentation

    prs = Presentation(io.BytesIO(file_bytes))
    text_parts = []

    for i, slide in enumerate(prs.slides):
        slide_texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_texts.append(shape.text.strip())
        if slide_texts:
            text_parts.append(f"--- Slide {i + 1} ---\n" + "\n".join(slide_texts))

    return "\n\n".join(text_parts) if text_parts else "No readable text found in slides."


def extract_text_from_docx(file_bytes: bytes) -> str:
    """Extract text from a Word document."""
    from docx import Document

    doc = Document(io.BytesIO(file_bytes))
    text_parts = []

    for para in doc.paragraphs:
        if para.text.strip():
            text_parts.append(para.text.strip())

    return "\n\n".join(text_parts) if text_parts else "No readable text found in document."


def extract_text_from_image(file_bytes: bytes, filename: str) -> str:
    """For images, return a note that the file is an image (AI will handle via description)."""
    return f"[Image file uploaded: {filename}. The AI should generate a general study schedule based on the subject and deadline provided.]"


def extract_text(file_bytes: bytes, filename: str) -> str:
    """
    Main extraction function - routes to the correct extractor based on file type.
    """
    ext = Path(filename).suffix.lower()

    extractors = {
        ".pdf": extract_text_from_pdf,
        ".pptx": extract_text_from_pptx,
        ".docx": extract_text_from_docx,
        ".png": lambda b: extract_text_from_image(b, filename),
        ".jpg": lambda b: extract_text_from_image(b, filename),
        ".jpeg": lambda b: extract_text_from_image(b, filename),
    }

    extractor = extractors.get(ext)
    if not extractor:
        return f"Unsupported file type: {ext}"

    try:
        text = extractor(file_bytes)
        # Limit text to ~15000 chars to stay within AI token limits
        if len(text) > 15000:
            text = text[:15000] + "\n\n[Content truncated for processing...]"
        return text
    except Exception as e:
        return f"Error extracting text from {filename}: {str(e)}"


def chunk_text_for_vectors(text: str, chunk_size: int = 900, overlap: int = 150) -> List[str]:
    """
    Split extracted text into overlapping chunks for vector storage/retrieval.
    """
    cleaned = (text or "").strip()
    if not cleaned:
        return []

    if chunk_size <= overlap:
        overlap = max(0, chunk_size // 5)

    chunks: List[str] = []
    start = 0
    length = len(cleaned)

    while start < length:
        end = min(length, start + chunk_size)
        chunk = cleaned[start:end].strip()
        if chunk:
            chunks.append(chunk)
        if end >= length:
            break
        start = max(0, end - overlap)

    return chunks